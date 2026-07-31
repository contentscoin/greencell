# -*- coding: utf-8 -*-
"""
GREENCELL BIOME-7 / LOVION
PDRN 강조 제품소개서 슬라이드 생성 스크립트
"""
import math
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from lxml import etree

# ---------------------------------------------------------------- palette
DARK      = RGBColor(0x07, 0x1D, 0x16)
DARK2     = RGBColor(0x0C, 0x2A, 0x20)
CARD_DARK = RGBColor(0x10, 0x38, 0x2A)
NEON      = RGBColor(0x2E, 0xE8, 0x7F)
NEON_DIM  = RGBColor(0x1C, 0x8F, 0x55)
MINT      = RGBColor(0x8F, 0xE3, 0xC4)
CREAM     = RGBColor(0xFD, 0xF9, 0xEE)
CARD_LT   = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x14, 0x2B, 0x22)
GRAY      = RGBColor(0x8F, 0xA5, 0x99)
GRAY_LT   = RGBColor(0x6B, 0x7D, 0x74)
GOLD      = RGBColor(0xC9, 0xB0, 0x74)
AMBER     = RGBColor(0xC8, 0x7A, 0x1E)
AMBER_BG  = RGBColor(0xFD, 0xF3, 0xDF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

KR = "맑은 고딕"
EN = "Segoe UI"
EN_L = "Segoe UI Light"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]

# ------------------------------------------------- 원본 PDF 추출 이미지 사용
from PIL import Image as PILImage

IMG = "/projects/sandbox/images"
DECK_IMG = f"{IMG}/_deck"
os.makedirs(DECK_IMG, exist_ok=True)


def _derive(src, out, box=None, ratio=None):
    """슬라이드 삽입용 파생 이미지 생성 (크롭)."""
    if os.path.exists(out):
        return out
    im = PILImage.open(src).convert("RGB")
    if box:
        im = im.crop(box)
    if ratio:                                  # 센터 크롭으로 비율 맞추기
        w, h = im.size
        th = int(w / ratio)
        top = max(0, (h - th) // 2)
        im = im.crop((0, top, w, min(h, top + th)))
    im.save(out)
    return out


COVER_BG = _derive(f"{IMG}/04_배경/p01_배경_숲반사_커버_951x585.png",
                   f"{DECK_IMG}/cover_bg_16x9.png", ratio=16 / 9)
BOTTLE_ROSE = _derive(f"{IMG}/01_제품사진/p24_제품_보틀_로즈박스_문구포함_2358x1651.png",
                      f"{DECK_IMG}/bottle_rose.png", box=(300, 250, 1150, 1290))
BOTTLE_MARBLE = f"{IMG}/01_제품사진/p02_제품_보틀_대리석_꽃_934x848.png"
BOTTLE_LAB = f"{IMG}/01_제품사진/p23_제품_보틀_실험실_1225x822.png"
RITUAL_ART = f"{IMG}/03_성분도판/p08_사용법_얼굴마사지_일러스트_814x1047.png"
PDRN_DIAGRAM = f"{IMG}/03_성분도판/p12_PDRN_작용기전_도식_885x375.png"


# ---------------------------------------------------------------- helpers
def set_typeface(run, name):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


def spacing(run, pts):
    """letter-spacing in points"""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def new_slide(bg=DARK):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def tb(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def put(tf, text, size, color, bold=False, font=KR, ls=None, first=False,
        space_before=0, space_after=0, align=None, spc=None, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    if ls:
        p.line_spacing = ls
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    set_typeface(r, font)
    if spc:
        spacing(r, spc)
    return p


def text(slide, x, y, w, h, content, size, color, bold=False, font=KR,
         ls=1.25, align=PP_ALIGN.LEFT, spc=None, anchor=MSO_ANCHOR.TOP,
         italic=False):
    """content: str or list[str]"""
    tf = tb(slide, x, y, w, h, align=align, anchor=anchor)
    lines = [content] if isinstance(content, str) else content
    for i, line in enumerate(lines):
        put(tf, line, size, color, bold=bold, font=font, ls=ls,
            first=(i == 0), align=align, spc=spc, italic=italic)
    return tf


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE,
         line=None, line_w=1.0, adj=None, alpha=None):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if color is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        if alpha is not None:
            sf = sh.fill.fore_color._xFill.find(qn('a:srgbClr'))
            a = etree.SubElement(sf, qn('a:alpha'))
            a.set('val', str(int(alpha * 100000)))
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    if adj is not None and sh.adjustments:
        sh.adjustments[0] = adj
    sh.text_frame.word_wrap = True
    return sh


def section_title(slide, en, kr=None, dark=True, y=0.62):
    """neon bar + english headline (+ korean sub headline)"""
    rect(slide, 0.72, y + 0.04, 0.075, 0.52, NEON)
    text(slide, 0.95, y, 9.5, 0.6,
         en, 27, WHITE if dark else INK, font=EN_L, spc=1.4)
    if kr:
        text(slide, 0.98, y + 0.66, 10.5, 0.45, kr, 14,
             MINT if dark else NEON_DIM, font=KR, spc=0.4)


PAGE_STYLE = {}   # id(slide) -> dark flag (번호는 최종 순서 확정 후 부여)


def page_num(slide, n=None, dark=True):
    """실제 번호는 stamp_page_numbers()에서 최종 슬라이드 순서 기준으로 찍는다."""
    PAGE_STYLE[id(slide)] = dark


def stamp_page_numbers():
    for i, slide in enumerate(prs.slides, start=1):
        if id(slide) not in PAGE_STYLE:
            continue          # 커버 등 번호 미표기 슬라이드
        dark = PAGE_STYLE[id(slide)]
        text(slide, SW - 1.25, SH - 0.62, 0.6, 0.3, f"{i:02d}", 11,
             GRAY if dark else GRAY_LT, font=EN, align=PP_ALIGN.RIGHT)


def brand_mark(slide, dark=True):
    tf = tb(slide, 0.72, SH - 0.66, 7.5, 0.3)
    r = tf.paragraphs[0].add_run()
    r.text = "LOVION"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = NEON if dark else NEON_DIM
    set_typeface(r, EN)
    spacing(r, 2.0)
    r2 = tf.paragraphs[0].add_run()
    r2.text = "   GREENCELL BIOME-7™  ALL-IN-ONE CREAM"
    r2.font.size = Pt(9)
    r2.font.color.rgb = GRAY if dark else GRAY_LT
    set_typeface(r2, EN)
    spacing(r2, 1.2)


def footnote(slide, msg, dark=True, y=None):
    y = SH - 0.98 if y is None else y
    text(slide, 0.72, y, 11.9, 0.34, msg, 9,
         GRAY if dark else GRAY_LT, font=KR, ls=1.15)


def helix(slide, cx, top, height, turns=2.4, dots=34, r1=0.085, r2=0.05,
          c1=NEON, c2=MINT, spread=0.95, alpha1=None):
    """decorative DNA double helix drawn with dots"""
    for i in range(dots):
        t = i / (dots - 1)
        ang = t * turns * 2 * math.pi
        y = top + t * height
        for phase, rr, cc in ((0, r1, c1), (math.pi, r2, c2)):
            x = cx + math.sin(ang + phase) * spread
            depth = (math.cos(ang + phase) + 1) / 2  # 0..1
            size = rr * (0.55 + 0.75 * depth)
            e = rect(slide, x - size / 2, y - size / 2, size, size, cc,
                     shape=MSO_SHAPE.OVAL,
                     alpha=0.25 + 0.75 * depth)
        # rungs
        if i % 4 == 0:
            x1 = cx + math.sin(ang) * spread
            x2 = cx + math.sin(ang + math.pi) * spread
            lo, hi = sorted([x1, x2])
            if hi - lo > 0.12:
                bar = rect(slide, lo, y - 0.012, hi - lo, 0.024, NEON_DIM,
                           alpha=0.5)


def pic(slide, path, x, y, w=None, h=None, frame=0.0, frame_color=NEON):
    """사진 삽입 (w 또는 h 하나만 주면 원본 비율 유지). frame>0 이면 네온 오프셋 프레임."""
    with PILImage.open(path) as im:
        iw, ih = im.size
    if w is None:
        w = h * iw / ih
    if h is None:
        h = w * ih / iw
    if frame:
        rect(slide, x + frame, y + frame, w, h, None, line=frame_color, line_w=1.25)
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return w, h


def chip(slide, x, y, w, h, label, color_bg, color_tx, size=10, bold=True,
         adj=0.5, font=KR):
    sh = rect(slide, x, y, w, h, color_bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
              adj=adj)
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color_tx
    set_typeface(r, font)
    return sh


# ================================================================ SLIDE 1
s = new_slide(DARK)
pic(s, COVER_BG, 0, 0, w=SW)                      # 원본 커버 배경 사진
rect(s, 0, 0, SW, SH, DARK, alpha=0.62)           # 텍스트 가독성 오버레이
rect(s, 0, 0, SW, 0.09, NEON, alpha=0.85)

text(s, 1.05, 1.25, 6.0, 0.32, "CLEAN SCIENTIFIC COSMETIC", 11, NEON,
     font=EN, spc=3.2, bold=True)

tf = tb(s, 1.0, 1.85, 8.6, 1.7)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "PDRN"
r.font.size = Pt(96); r.font.bold = True; r.font.color.rgb = NEON
set_typeface(r, EN); spacing(r, 2.0)
r = p.add_run(); r.text = "  ×"
r.font.size = Pt(40); r.font.color.rgb = MINT
set_typeface(r, EN)

text(s, 1.05, 3.35, 9.0, 0.75, "GREENCELL BIOME-7™", 38, WHITE,
     font=EN_L, spc=1.6)
rect(s, 1.08, 4.22, 2.6, 0.055, NEON)

text(s, 1.05, 4.62, 8.2, 1.2,
     ["DNA 유래 재생 과학을 담은 5세대 올인원 크림",
      "PDRN CORE + 6 POWER COMPLEX, 한 병으로 완성하는 고농축 스킨케어"],
     16, MINT, ls=1.5)

tf = tb(s, 1.05, 6.15, 8.0, 0.4)
r = tf.paragraphs[0].add_run()
r.text = "LOVION"
r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
set_typeface(r, EN); spacing(r, 3.0)
r = tf.paragraphs[0].add_run()
r.text = "      Product Introduction  ·  Shelly for U"
r.font.size = Pt(11); r.font.color.rgb = GRAY
set_typeface(r, EN)

pic(s, BOTTLE_MARBLE, 8.62, 1.95, w=3.95, frame=0.16)   # 제품 사진

footnote(s, "※ PDRN 전성분 표기명: 소듐디엔에이(Sodium DNA). 함량·유래는 제조사 확인 후 반영 예정.")

# ================================================================ SLIDE 2
s = new_slide(CREAM)
section_title(s, "Why PDRN Now", "피부 재생 성분 시장의 중심, PDRN", dark=False)
rect(s, 0, 0, 0.28, SH, NEON_DIM)

cards = [
    ("MARKET", "시술에서 홈케어로",
     ["재생·리페어 카테고리가 뷰티 시장의",
      "핵심 축으로 이동",
      "‘시술 후 홈케어’ 수요 지속 확대"]),
    ("CONSUMER", "성분을 검증하는 소비자",
     ["성분표를 직접 확인하고 비교하는",
      "성분 리터러시 세대",
      "‘무엇이 들었는지’가 구매 기준"]),
    ("BRAND", "과학으로 증명하는 브랜드",
     ["제약 기반 배합 기술력 +",
      "PDRN CORE 설계",
      "클린 사이언티픽 코스메틱 포지셔닝"]),
]
x = 0.72
for en, kr, lines in cards:
    c = rect(s, x, 2.15, 3.82, 2.75, CARD_LT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.05)
    rect(s, x, 2.15, 3.82, 0.075, NEON_DIM)
    text(s, x + 0.42, 2.52, 3.0, 0.3, en, 10, NEON_DIM, font=EN, spc=2.4, bold=True)
    text(s, x + 0.42, 2.92, 3.1, 0.5, kr, 19, INK, bold=True)
    text(s, x + 0.42, 3.62, 3.1, 1.2, lines, 12, GRAY_LT, ls=1.5)
    x += 4.03

bar = rect(s, 0.72, 5.28, 11.9, 1.05, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.12)
text(s, 1.15, 5.52, 3.4, 0.6, "KEY MESSAGE", 10, NEON_DIM, font=EN, spc=2.2, bold=True)
text(s, 1.15, 5.82, 10.9, 0.5,
     "“재생 성분 PDRN을 중심에 두고, 장벽·수분·마이크로바이옴까지 한 번에 잡는 올인원 설계.”",
     16, INK, bold=True)

footnote(s, "※ 시장 규모·성장률 등 수치 데이터는 확정된 출처 확인 후 삽입 예정 [DATA TBD]", dark=False)
brand_mark(s, dark=False)
page_num(s, 2, dark=False)

# ================================================================ SLIDE 3
s = new_slide(DARK)
rect(s, 6.6, 0, 6.733, SH, DARK2)
section_title(s, "What is PDRN", "폴리데옥시리보뉴클레오타이드 — 자연 유래 DNA 단편 복합체")

pic(s, BOTTLE_LAB, 7.35, 2.35, w=5.1, frame=0.16)       # 실험실 제품 사진
text(s, 7.35, 6.04, 5.1, 0.3, "GREENCELL BIOME-7™  ·  Formulation Lab", 9.5,
     GRAY, font=EN, spc=1.4)

text(s, 0.95, 2.15, 5.3, 1.6,
     ["PDRN은 살아 있는 세포에서 추출한",
      "저분자 DNA 단편의 혼합물로,",
      "피부가 스스로 회복하는 과정에",
      "필요한 뉴클레오타이드를 공급합니다."],
     15, WHITE, ls=1.55)

specs = [
    ("정의", "Polydeoxyribonucleotide", "저분자 DNA 단편 복합체"),
    ("구성", "Nucleotide / Nucleoside", "아데노신 등 활성 단위 포함"),
    ("유래", "[확인 중]", "원료 규격서 확인 후 확정"),
    ("특징", "저자극 · 고안정 뉴클레오타이드", "장기 사용에 적합한 프로파일"),
]
y = 3.82
for k, v, sub in specs:
    rect(s, 0.95, y + 0.05, 0.055, 0.5, NEON)
    text(s, 1.2, y, 0.9, 0.3, k, 11, MINT, bold=True)
    text(s, 2.15, y - 0.03, 4.1, 0.32, v, 13, WHITE, bold=True)
    text(s, 2.15, y + 0.27, 4.1, 0.3, sub, 9.5, GRAY)
    y += 0.66

footnote(s, "※ PDRN의 전성분 표기명은 소듐디엔에이(Sodium DNA)입니다.")
brand_mark(s)
page_num(s, 3)

# ================================================================ SLIDE 4
s = new_slide(DARK)
section_title(s, "How PDRN Works", "A2A 수용체를 통한 4가지 작용 메커니즘 (원료 학술자료 기준)")

# 원료사 기전 도판 (원본 소개서 p12)
rect(s, 0.62, 2.12, 7.94, 3.62, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.05)
pic(s, PDRN_DIAGRAM, 0.78, 2.42, w=7.6)
text(s, 0.78, 5.85, 7.6, 0.3, "PDRN 작용 기전 도판 — 원료 학술자료 기준", 9.5, GRAY)

text(s, 8.95, 1.92, 4.0, 0.3, "4 KEY MECHANISMS", 9.5, NEON, font=EN,
     spc=2.2, bold=True)

mechs = [
    ("01", "항염 밸런싱",
     "TNF-α·IL-6 신호는 낮추고 IL-10 활성은 높이는 방향"),
    ("02", "혈관 확장 · 순환",
     "VEGF 분비 촉진으로 미세순환·영양 공급 경로 활성"),
    ("03", "성장인자 시너지",
     "EGF·FGF·IGF 분비를 도와 피부 세포 활동 뒷받침"),
    ("04", "DNA 합성 살베이지",
     "우회 경로 활성화로 적은 에너지로 세포 재생 유도"),
]
y = 2.38
for no, title, desc in mechs:
    rect(s, 8.75, y + 0.04, 0.05, 0.66, NEON)
    text(s, 8.95, y - 0.02, 0.5, 0.3, no, 13, NEON, font=EN, bold=True)
    text(s, 9.42, y, 3.2, 0.3, title, 13.5, WHITE, bold=True)
    text(s, 8.95, y + 0.36, 3.65, 0.45, desc, 10, MINT, ls=1.35)
    y += 0.86

footnote(s, "※ 상기 내용은 PDRN 원료에 대한 일반 학술자료를 요약한 설명으로, 완제품의 의학적 효능·효과를 의미하지 않습니다.")
brand_mark(s)
page_num(s, 4)

# ================================================================ SLIDE 5
s = new_slide(CREAM)
section_title(s, "Expected Skin Benefits", "PDRN CORE 설계가 향하는 피부 변화", dark=False)
rect(s, 0, 0, 0.28, SH, NEON_DIM)

items = [
    ("진정", "예민해진 피부를 편안하게", "자극 후 붉은 기와 당김을 진정시키는 저자극 설계"),
    ("장벽", "속부터 단단해지는 피부", "세라마이드 복합체와 함께 무너진 장벽을 채워줌"),
    ("탄력", "밀도가 느껴지는 피부결", "5GF 나노좀과의 조합으로 탄탄한 볼륨감 서포트"),
    ("수분", "속건조 없는 촉촉함", "8중 히알루론산으로 층층이 수분 레이어링"),
    ("밸런스", "건강한 피부 환경 유지", "마이크로바이옴 3종으로 유익균 밸런스 관리"),
]
y = 2.12
for i, (k, t, d) in enumerate(items):
    row = rect(s, 0.72, y, 11.9, 0.78, CARD_LT if i % 2 == 0 else RGBColor(0xF7, 0xF3, 0xE6),
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.18)
    chip(s, 0.95, y + 0.16, 1.05, 0.46, k, NEON_DIM, WHITE, size=11.5)
    text(s, 2.25, y + 0.19, 4.2, 0.4, t, 15, INK, bold=True)
    text(s, 6.55, y + 0.23, 5.9, 0.4, d, 12, GRAY_LT)
    y += 0.9

footnote(s, "※ 화장품 표시·광고 범위 내 표현입니다. 의약품적 효능(치료·재생 등) 표현은 사용하지 않습니다.", dark=False)
brand_mark(s, dark=False)
page_num(s, 5)

# ================================================================ SLIDE 6
s = new_slide(DARK)
rect(s, 0, 0, SW, SH, DARK2, alpha=0.4)
section_title(s, "PDRN CORE + 6 Complex", "PDRN을 중심에 둔 골든 레이시오 배합 구조")

# center core
core = rect(s, 5.14, 2.55, 3.05, 3.05, NEON_DIM, shape=MSO_SHAPE.OVAL)
rect(s, 4.84, 2.25, 3.65, 3.65, NEON, shape=MSO_SHAPE.OVAL, alpha=0.12)
tf = core.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "PDRN"
r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE
set_typeface(r, EN); spacing(r, 1.5)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "CORE"
r.font.size = Pt(12); r.font.color.rgb = MINT
set_typeface(r, EN); spacing(r, 3.0)
text(s, 4.4, 5.70, 4.55, 0.32, "재생 시그널의 중심", 11.5, WHITE,
     align=PP_ALIGN.CENTER, bold=True)

left = [
    ("5GF NANOSOME", "성장인자 5종", "sh-Oligopeptide-1·2 / sh-Polypeptide-1·3·16"),
    ("8-FOLD HA", "8중 히알루론산", "저분자~고분자 수분 레이어링"),
    ("CERAMIDE", "복합 세라마이드 NP", "AS·NS·AP·EOP·NG·EOS 장벽 강화"),
]
right = [
    ("MICROBIOME", "마이크로바이옴 3종", "비피다·락토바실러스·스트렙토코쿠스 발효"),
    ("PLANT STEM CELL", "식물줄기세포 배양액 5종", "캘러스 배양 추출물 기반 활력 부여"),
    ("PEPTIDE", "펩타이드 5종", "피부 구성 단백질 유사 구조로 탄력 서포트"),
]
y0 = 2.22
for i, (en, kr, sub) in enumerate(left):
    y = y0 + i * 1.32
    rect(s, 0.72, y, 4.05, 1.16, CARD_DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.09)
    rect(s, 0.72, y + 0.12, 0.055, 0.92, NEON)
    text(s, 0.98, y + 0.16, 3.6, 0.28, en, 9, NEON, font=EN, spc=1.8, bold=True)
    text(s, 0.98, y + 0.44, 3.6, 0.3, kr, 14, WHITE, bold=True)
    text(s, 0.98, y + 0.78, 3.7, 0.3, sub, 9.5, GRAY)
    rect(s, 4.82, y + 0.55, 0.34, 0.03, NEON, alpha=0.5)

for i, (en, kr, sub) in enumerate(right):
    y = y0 + i * 1.32
    rect(s, 8.56, y, 4.05, 1.16, CARD_DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.09)
    rect(s, 12.555, y + 0.12, 0.055, 0.92, NEON)
    text(s, 8.82, y + 0.16, 3.6, 0.28, en, 9, NEON, font=EN, spc=1.8, bold=True)
    text(s, 8.82, y + 0.44, 3.6, 0.3, kr, 14, WHITE, bold=True)
    text(s, 8.82, y + 0.78, 3.7, 0.3, sub, 9.5, GRAY)
    rect(s, 8.17, y + 0.55, 0.34, 0.03, NEON, alpha=0.5)

text(s, 4.4, 6.14, 4.55, 0.35, "GOLDEN RATIO FORMULATION", 10, MINT, font=EN,
     spc=2.6, align=PP_ALIGN.CENTER, bold=True)

footnote(s, "※ PDRN 전성분 표기명: 소듐디엔에이(Sodium DNA). 함량은 제조사 확인 후 표기 예정.")
brand_mark(s)
page_num(s, 6)

# ================================================================ SLIDE 7
s = new_slide(CREAM)
section_title(s, "Synergy Logic", "PDRN이 켜고, 6가지 콤플렉스가 완성하는 4단계 케어", dark=False)
rect(s, 0, 0, 0.28, SH, NEON_DIM)

steps = [
    ("STEP 01", "SIGNAL", "PDRN CORE",
     ["재생 시그널을 깨우는", "뉴클레오타이드 공급"]),
    ("STEP 02", "AMPLIFY", "5GF · 펩타이드",
     ["성장인자·펩타이드로", "탄력과 밀도 신호 증폭"]),
    ("STEP 03", "LOCK", "세라마이드 · 8중 HA",
     ["장벽을 세우고 수분을", "층층이 고정"]),
    ("STEP 04", "BALANCE", "마이크로바이옴",
     ["유익균 밸런스로", "건강한 피부 환경 유지"]),
]
x = 0.72
for i, (st, en, kr, lines) in enumerate(steps):
    w = 2.78
    card = rect(s, x, 2.2, w, 2.85, CARD_LT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.07)
    rect(s, x, 2.2, w, 0.075, NEON_DIM)
    text(s, x + 0.32, 2.55, 2.0, 0.28, st, 9.5, NEON_DIM, font=EN, spc=2.0, bold=True)
    text(s, x + 0.32, 2.88, 2.2, 0.45, en, 21, INK, font=EN_L, spc=0.8)
    rect(s, x + 0.34, 3.48, 0.85, 0.045, NEON_DIM)
    text(s, x + 0.32, 3.72, 2.3, 0.35, kr, 13.5, NEON_DIM, bold=True)
    text(s, x + 0.32, 4.18, 2.3, 0.7, lines, 11.5, GRAY_LT, ls=1.45)
    if i < 3:
        ar = rect(s, x + w + 0.055, 3.42, 0.36, 0.4, NEON_DIM,
                  shape=MSO_SHAPE.CHEVRON, alpha=0.85)
    x += w + 0.47

bar = rect(s, 0.72, 5.45, 11.9, 0.95, RGBColor(0x10, 0x38, 0x2A),
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.14)
text(s, 1.15, 5.72, 11.2, 0.5,
     "한 병으로 끝나는 재생 루틴 — 복잡한 단계 없이, 올인원 크림 하나로 4단계 케어를 완성합니다.",
     15.5, WHITE, bold=True)
brand_mark(s, dark=False)
page_num(s, 7)

# ================================================================ SLIDE 8
s = new_slide(DARK)
rect(s, 6.85, 0, 6.5, SH, DARK2)
section_title(s, "Product & Ritual", "제품 정보와 PDRN 흡수를 높이는 사용 리추얼")

specs = [
    ("제품명", "그린셀바이옴 올인원 크림 (GREENCELL BIOME-7™ All-in-One Cream)"),
    ("브랜드", "LOVION  ·  Clean Scientific Cosmetic"),
    ("용량", "50ml (1.69 fl. oz.)"),
    ("핵심 설계", "PDRN CORE + 5GF 나노좀 + 8중 HA + 복합 세라마이드 + 마이크로바이옴"),
    ("사용 대상", "젠더리스 · 전 연령 데일리 케어"),
    ("사용 단계", "세안 후 마지막 단계, 1일 2회 (아침 · 저녁)"),
]
y = 2.18
for k, v in specs:
    text(s, 0.95, y, 1.5, 0.3, k, 11, NEON, bold=True)
    text(s, 2.42, y - 0.02, 4.1, 0.55, v, 11.5, WHITE, ls=1.35)
    rect(s, 0.95, y + 0.56, 5.55, 0.012, RGBColor(0x1E, 0x4A, 0x38))
    y += 0.72

text(s, 7.35, 2.18, 5.2, 0.32, "LAYERING RITUAL", 10, NEON, font=EN, spc=2.6, bold=True)

# 사용법 일러스트 (원본 소개서 p08) — 흰 배경이므로 카드로 감싼다
rect(s, 10.55, 2.45, 2.34, 2.94, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.06)
pic(s, RITUAL_ART, 10.67, 2.57, h=2.70)

ritual = [
    ("1", "Neck Line", ["목선부터 아래에서 위로", "쓸어 올리며 가볍게 도포"]),
    ("2", "Jaw & Cheek", ["턱선·볼·눈가를 리프팅", "방향으로 꼼꼼히 도포"]),
    ("3", "Forehead", ["눈썹을 자극하며", "이마까지 마무리"]),
]
y = 2.62
for no, en, kr in ritual:
    circ = rect(s, 7.35, y, 0.55, 0.55, NEON_DIM, shape=MSO_SHAPE.OVAL)
    tfc = circ.text_frame
    tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
    pc = tfc.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
    rc = pc.add_run(); rc.text = no
    rc.font.size = Pt(14); rc.font.bold = True; rc.font.color.rgb = WHITE
    set_typeface(rc, EN)
    text(s, 8.08, y - 0.01, 2.4, 0.3, en, 13.5, WHITE, font=EN, bold=True)
    text(s, 8.08, y + 0.32, 2.5, 0.6, kr, 10.5, MINT, ls=1.35)
    y += 1.02

tip = rect(s, 7.35, 5.72, 5.55, 0.68, NEON, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
           adj=0.18, alpha=0.14)
text(s, 7.62, 5.87, 5.1, 0.45, "TIP  ·  2~3번 덧바르면 더 깊은 보습과 탄력을 경험할 수 있습니다.",
     11, MINT, bold=True)
footnote(s, "※ 용량은 패키지 라벨 표기(50ml / 1.69 fl. oz.) 기준입니다. 제조사·책임판매업자 등 나머지 표시사항은 최종 패키지 기준으로 확정합니다.")
brand_mark(s)
page_num(s, 8)

# ================================================================ SLIDE 9
s = new_slide(WHITE)
rect(s, 0, 0, 0.28, SH, AMBER)
section_title(s, "Labeling & Claim Check", "PDRN 강조 전 반드시 확인해야 하는 항목", dark=False)

warn = rect(s, 0.72, 2.02, 11.9, 1.18, AMBER_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.1)
rect(s, 0.72, 2.02, 0.09, 1.18, AMBER)
text(s, 1.05, 2.24, 11.3, 0.35, "표기명 확정 · 나머지 항목 확인 중", 12, AMBER, bold=True)
text(s, 1.05, 2.60, 11.3, 0.55,
     ["PDRN의 전성분 표기명은 소듐디엔에이(Sodium DNA)로 확정되었습니다.",
      "함량·유래 등 나머지 항목은 제조사 확인 자료 수령 후 반영합니다."],
     11, INK, ls=1.4)

text(s, 0.95, 3.55, 5.6, 0.3, "제조사 확인 요청 항목", 13, INK, bold=True)
asks = [
    "PDRN 원료의 실제 배합 여부",
    "전성분표상 정확한 표기 성분명",
    "완제품 내 실제 함량 (ppm 또는 %)",
    "유래 (연어 / 식물 등) 및 원산지",
    "원료 규격서 · 시험성적서",
    "기존 생산분 / 신규 처방 적용 범위",
]
y = 3.92
for a in asks:
    chip(s, 0.95, y + 0.02, 0.22, 0.22, "", NEON_DIM, WHITE, size=8, adj=0.5)
    text(s, 1.38, y - 0.01, 5.0, 0.3, a, 11.5, GRAY_LT)
    y += 0.40

text(s, 7.0, 3.55, 5.6, 0.3, "표현 가이드", 13, INK, bold=True)
ok = rect(s, 7.0, 3.95, 5.62, 1.05, RGBColor(0xEF, 0xFA, 0xF3),
          shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.12)
text(s, 7.28, 4.10, 5.1, 0.3, "사용 가능 (확정 후)", 10.5, NEON_DIM, bold=True)
text(s, 7.28, 4.42, 5.1, 0.5,
     "“PDRN 성분 함유”, “피부 진정·탄력에 도움”, “장벽 강화 보습”", 11, INK, ls=1.35)

no = rect(s, 7.0, 5.18, 5.62, 1.05, RGBColor(0xFD, 0xF0, 0xEF),
          shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.12)
text(s, 7.28, 5.33, 5.1, 0.3, "사용 불가", 10.5, RGBColor(0xC0, 0x39, 0x2B), bold=True)
text(s, 7.28, 5.65, 5.1, 0.5,
     "“세포 재생·치료·주사 효과 동일”, “의약품적 효능” 표현", 11, INK, ls=1.35)

text(s, 0.95, 6.42, 5.9, 0.35,
     "※ 발효추출물·캘러스배양추출물은 PDRN과 동일 원료로 볼 수 없습니다.", 10, AMBER, bold=True)
brand_mark(s, dark=False)
page_num(s, 9, dark=False)

# ================================================================ SLIDE 10
s = new_slide(DARK)
rect(s, 0, 0, SW, SH, DARK2, alpha=0.5)
rect(s, 0, SH - 0.09, SW, 0.09, NEON)
_w, _h = pic(s, BOTTLE_ROSE, 0.85, 1.30, h=4.9)         # 로즈박스 제품 사진
rect(s, 0.85, 1.30, _w, _h, None, line=NEON_DIM, line_w=1.0)

text(s, 5.1, 2.35, 7.4, 0.35, "PDRN × GREENCELL BIOME-7™", 12, NEON, font=EN,
     spc=2.6, bold=True)
text(s, 5.1, 2.85, 7.5, 2.3,
     ["건강하게", "아름다워지는", "그날까지…"], 40, WHITE, font=EN_L, ls=1.22)
rect(s, 5.14, 5.32, 2.2, 0.05, NEON)
text(s, 5.1, 5.62, 7.4, 0.9,
     ["Clean Scientific Cosmetic", "LOVION  ·  Shelly for U"],
     14, MINT, font=EN, ls=1.5)
page_num(s, 10)

# ================================================ SLIDE A (비교) → 최종 05
s = new_slide(CREAM)
rect(s, 0, 0, 0.28, SH, NEON_DIM)
section_title(s, "Ingredient Positioning",
              "PDRN과 유사 재생 성분 — 역할이 어떻게 다른가", dark=False)

COLS = [(0.72, 2.35), (3.13, 3.05), (6.24, 3.40), (9.70, 2.92)]
HEADS = ["성분", "핵심 역할", "작용 방식", "포지션"]
ROWS = [
    ("PDRN", "뉴클레오타이드 공급",
     "DNA 단편이 재생에 쓰이는 원재료를 직접 공급",
     "재생 설계의 출발점", "core"),
    ("성장인자 (5GF)", "세포 신호 전달",
     "성장인자가 피부 구성 세포의 활동을 지시",
     "PDRN과 상호 보완 관계", "n"),
    ("8중 히알루론산", "수분 저장",
     "분자량별로 겉·속 수분을 채우고 붙잡아 둠",
     "재생이 아닌 ‘보습’ 축", "n"),
    ("복합 세라마이드", "장벽 구성",
     "장벽 지질을 채워 수분 손실을 물리적으로 차단",
     "재생이 아닌 ‘방어’ 축", "n"),
    ("발효 · 캘러스 추출물", "진정 · 컨디셔닝",
     "발효 유효성분·식물 배양 추출물의 컨디셔닝",
     "PDRN 대체 원료가 아님", "warn"),
]

for (cx, cw), h in zip(COLS, HEADS):
    hd = rect(s, cx, 1.94, cw, 0.44, NEON_DIM, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.25)
    text(s, cx + 0.22, 2.05, cw - 0.4, 0.3, h, 11, WHITE, bold=True)

y = 2.48
for name, role, how, pos, kind in ROWS:
    if kind == "core":
        bg, c_main, c_sub, c_pos = CARD_DARK, WHITE, MINT, NEON
    elif kind == "warn":
        bg, c_main, c_sub, c_pos = AMBER_BG, INK, GRAY_LT, AMBER
    else:
        bg, c_main, c_sub, c_pos = CARD_LT, INK, GRAY_LT, NEON_DIM
    for cx, cw in COLS:
        rect(s, cx, y, cw, 0.70, bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.14)
    if kind == "core":
        rect(s, COLS[0][0], y, 0.06, 0.70, NEON)
    text(s, COLS[0][0] + 0.24, y + 0.21, COLS[0][1] - 0.42, 0.4, name, 12.5,
         c_main, bold=True)
    text(s, COLS[1][0] + 0.22, y + 0.23, COLS[1][1] - 0.4, 0.4, role, 11.5, c_main)
    text(s, COLS[2][0] + 0.22, y + 0.24, COLS[2][1] - 0.4, 0.4, how, 10.5, c_sub)
    text(s, COLS[3][0] + 0.22, y + 0.24, COLS[3][1] - 0.4, 0.4, pos, 10.5,
         c_pos, bold=True)
    y += 0.78

footnote(s, "※ 성분별 역할을 비교한 일반 정보이며 제품 간 효능 비교가 아닙니다. PDRN의 배합·표기명·함량은 제조사 확인 후 확정합니다.",
         dark=False)
brand_mark(s, dark=False)
page_num(s, dark=False)

# ================================================ SLIDE B (FAQ) → 최종 10
s = new_slide(DARK)
rect(s, 0, 0, SW, SH, DARK2, alpha=0.35)
section_title(s, "FAQ & Sales Guide", "현장에서 자주 받는 질문과 표준 응대")

faqs = [
    (0.72, 2.15, "Q1", "PDRN이 실제로 들어 있나요?",
     "전성분표에 소듐디엔에이(Sodium DNA)로 표기된 성분이\nPDRN입니다. 함량은 확인 후 안내드립니다."),
    (6.85, 2.15, "Q2", "시술(주사)과 같은 효과인가요?",
     "아닙니다. 화장품은 피부 표면 케어를 목적으로 하며,\n의학적 시술과 동일한 효과로 설명할 수 없습니다."),
    (0.72, 4.32, "Q3", "유래가 어떻게 되나요? 비건인가요?",
     "유래(연어/식물 등)는 원료 규격서 확인 중입니다.\n확정되면 안내드리겠습니다."),
    (6.85, 4.32, "Q4", "민감성 피부도 사용할 수 있나요?",
     "저자극 설계이지만, 예민한 피부는 팔 안쪽에\n패치 테스트 후 사용을 권장합니다."),
]
for x, y, q, question, answer in faqs:
    rect(s, x, y, 5.76, 1.9, CARD_DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.08)
    rect(s, x, y, 5.76, 0.06, NEON)
    chip(s, x + 0.34, y + 0.3, 0.62, 0.34, q, NEON_DIM, WHITE, size=11, font=EN)
    text(s, x + 1.12, y + 0.34, 4.4, 0.35, question, 14, WHITE, bold=True)
    text(s, x + 0.34, y + 0.98, 5.1, 0.8, answer.split("\n"), 11, MINT, ls=1.45)

footnote(s, "※ 확정되지 않은 항목은 “확인 중”으로 안내하고 추정 답변을 하지 않습니다. 의약품적 효능·시술 대체 표현은 사용 금지입니다.")
brand_mark(s)
page_num(s)

# ------------------------------------------------------- 최종 슬라이드 순서
# 생성순: 0커버 1why 2what 3how 4benefit 5core 6synergy 7product 8labeling
#         9closing 10비교 11FAQ
ORDER = [0, 1, 2, 3, 10, 4, 5, 6, 7, 11, 8, 9]
sldIdLst = prs.slides._sldIdLst
_ids = list(sldIdLst)
for i in ORDER:
    sldIdLst.append(_ids[i])   # 기존 자식을 append하면 '이동'이 된다
stamp_page_numbers()

# ------------------------------------------------------- 발표용 스크립트(노트)
NOTES = [
    # 01 커버
    "[오프닝 30초]\n"
    "그린셀바이옴 올인원 크림을 ‘PDRN 중심 재생 케어’ 관점에서 소개하는 자료입니다.\n"
    "오늘 전달할 한 문장: “재생 성분 PDRN을 중심에 두고, 장벽·수분·마이크로바이옴까지 "
    "한 병으로 잡는 올인원 설계.”\n"
    "※ 내부 주의: PDRN 배합·표기명·함량은 제조사 확인 전입니다. "
    "확정 전까지 ‘PDRN 함유’를 단정적으로 말하지 말고 “PDRN 설계 방향”으로 설명하세요.",

    # 02 Why PDRN Now
    "[시장 배경 1분]\n"
    "재생·리페어가 뷰티 카테고리의 중심으로 이동했고, 시술 후 홈케어 수요가 계속 늘고 있습니다.\n"
    "요즘 고객은 전성분표를 직접 확인하고 비교합니다. 그래서 ‘무엇이 들었는지’가 구매 기준이 됩니다.\n"
    "우리는 제약 기반 배합 기술력 + PDRN CORE 설계로 이 흐름에 답합니다.\n"
    "※ 시장 수치는 출처 확정 후 삽입 예정이므로, 지금은 숫자를 말하지 마세요.",

    # 03 What is PDRN
    "[성분 정의 1분]\n"
    "PDRN은 폴리데옥시리보뉴클레오타이드, 쉽게 말해 저분자 DNA 단편의 복합체입니다.\n"
    "피부가 스스로 회복하는 과정에 필요한 ‘재료(뉴클레오타이드)’를 공급한다는 개념으로 설명하면 "
    "고객이 가장 쉽게 이해합니다.\n"
    "유래(연어/식물)와 전성분 표기명은 원료 규격서 확인 후 확정입니다. 질문받으면 “확인 중”으로 답하세요.",

    # 04 How PDRN Works
    "[기전 1분 30초]\n"
    "핵심은 A2A 수용체입니다. 여기서 네 가지 방향으로 작용한다고 보고되어 있습니다.\n"
    "① 염증 신호는 낮추고 항염 신호는 높이는 밸런싱, ② VEGF를 통한 순환·영양 공급, "
    "③ EGF·FGF·IGF 등 성장인자 시너지, ④ DNA 합성 살베이지 경로로 적은 에너지로 세포 재생 유도.\n"
    "※ 이 내용은 ‘원료에 대한 학술자료 요약’입니다. 완제품의 의학적 효능으로 말하면 광고 위반이 됩니다. "
    "반드시 “원료 연구에서는 이렇게 보고됩니다”라는 표현을 쓰세요.",

    # 05 Ingredient Positioning
    "[차별점 1분]\n"
    "고객이 가장 많이 혼동하는 부분입니다. 성분마다 역할이 다릅니다.\n"
    "PDRN은 재생의 ‘출발점(재료 공급)’, 성장인자는 ‘신호’, 히알루론산은 ‘보습’, 세라마이드는 ‘방어’입니다.\n"
    "특히 마지막 줄이 중요합니다. 발효추출물이나 캘러스배양추출물은 PDRN의 대체 원료가 아닙니다. "
    "‘식물 줄기세포니까 PDRN과 같다’는 설명은 절대 하지 마세요.",

    # 06 Expected Skin Benefits
    "[기대 변화 1분]\n"
    "진정 → 장벽 → 탄력 → 수분 → 밸런스, 다섯 가지 축으로 설명합니다.\n"
    "고객 상황에 맞춰 하나만 골라 강조하는 것이 효과적입니다. "
    "예: 시술 후에는 진정, 겨울철에는 장벽·수분, 40대 이상은 탄력·밀도.\n"
    "※ ‘재생’, ‘치료’ 같은 의약품적 표현 대신 “도움을 줍니다”, “관리합니다”로 말하세요.",

    # 07 PDRN CORE + 6 Complex
    "[배합 구조 1분]\n"
    "PDRN을 중심에 두고 6가지 콤플렉스가 둘러싸는 구조입니다.\n"
    "5GF 나노좀, 8중 히알루론산, 복합 세라마이드, 마이크로바이옴 3종, 식물줄기세포 배양액 5종, 펩타이드 5종.\n"
    "‘성분을 많이 넣었다’가 아니라 ‘역할이 겹치지 않게 배치했다’는 점을 강조하세요.\n"
    "※ 함량(ppm/%)은 제조사 확인 후 표기 예정입니다.",

    # 08 Synergy Logic
    "[스토리 마무리 1분]\n"
    "네 단계로 정리합니다. PDRN이 신호를 켜고(SIGNAL), 5GF·펩타이드가 증폭하고(AMPLIFY), "
    "세라마이드·히알루론산이 장벽과 수분을 고정하고(LOCK), 마이크로바이옴이 환경을 유지합니다(BALANCE).\n"
    "여기서 올인원의 가치가 나옵니다. 네 단계를 각각 다른 제품으로 사는 대신, 한 병으로 끝냅니다.\n"
    "가격 저항이 있는 고객에게는 ‘라인 구매 대비 비용’ 관점으로 연결하세요.",

    # 09 Product & Ritual
    "[제품 정보 + 시연 1분 30초]\n"
    "제품 스펙을 간단히 짚고, 오른쪽 사용 리추얼은 말로만 하지 말고 직접 손으로 시연하세요.\n"
    "목선 → 턱·볼·눈가 → 이마 순서로, 아래에서 위로 쓸어 올리는 방향이 핵심입니다.\n"
    "탄력이 떨어진 쪽을 먼저 레이어링하고 반대쪽도 같은 방식으로 진행하도록 안내합니다.\n"
    "2~3회 덧바르면 보습·탄력 체감이 커진다는 팁으로 마무리하세요.\n"
    "※ 용량 등 표시사항은 최종 패키지 기준으로 확정합니다.",

    # 10 FAQ & Sales Guide
    "[응대 가이드 — 내부 교육용]\n"
    "네 가지 질문은 현장에서 반드시 나옵니다. 답변 톤을 통일하세요.\n"
    "Q1은 지금 가장 민감합니다. 확정 자료가 없으면 “제조사 확인 자료 기준으로 안내드리겠습니다”로 "
    "끊고, 추정 답변을 하지 마세요.\n"
    "Q2는 반드시 선을 그어야 합니다. 시술과 동일한 효과라는 뉘앙스도 위험합니다.\n"
    "Q3·Q4는 확인 중 항목과 패치 테스트 권장으로 안내합니다.",

    # 11 Labeling & Claim Check
    "[내부 검토용 — 대외 발표 시 제외]\n"
    "현재 패키지 전성분표에서 PDRN에 해당하는 표기 성분명이 확인되지 않았습니다.\n"
    "따라서 확정 자료를 받기 전까지 PDRN 함유 표현은 대외 사용을 보류합니다.\n"
    "제조사에는 배합 여부·표기명·함량·유래·규격서·적용 범위 6개 항목을 문서로 요청하고, "
    "회신 자료를 근거 파일로 보관하세요.\n"
    "회신 후에는 왼쪽 초록 박스의 표현만 사용하고, 빨간 박스 표현은 어떤 채널에서도 쓰지 않습니다.",

    # 12 클로징
    "[클로징 30초]\n"
    "“건강하게 아름다워지는 그날까지” — 브랜드 메시지로 마무리합니다.\n"
    "마지막에 다음 액션을 반드시 제안하세요. 샘플 제공, 체험 일정, 또는 발주 상담 중 하나를 고르게 합니다.\n"
    "질문이 남았다면 확정되지 않은 항목은 메모해서 확인 후 회신하겠다고 안내하세요.",
]

for _slide, _note in zip(prs.slides, NOTES):
    _tf = _slide.notes_slide.notes_text_frame
    _tf.text = _note

# --------------------------------------------------------------- cleanup
# 마스터/레이아웃의 날짜·꼬리말·페이지번호 자동 placeholder 제거
from pptx.enum.shapes import PP_PLACEHOLDER
KILL = (PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER)
for master in prs.slide_masters:
    targets = [master] + list(master.slide_layouts)
    for t in targets:
        for shp in list(t.shapes):
            try:
                if shp.is_placeholder and shp.placeholder_format.type in KILL:
                    shp._element.getparent().remove(shp._element)
            except Exception:
                pass

out = "/projects/sandbox/GREENCELL_BIOME_PDRN_제품소개서.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
